#!/usr/bin/env python3
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


PRESENTATION_DIR = Path(__file__).resolve().parents[1]
ASSETS_DIR = PRESENTATION_DIR / "assets"
DIAGRAMS_DIR = PRESENTATION_DIR / "diagrams" / "rendered"
OUTPUT_DIR = PRESENTATION_DIR / "output"
OUTPUT_FILE = OUTPUT_DIR / "IMPORTANTE_MANUAL_defensa-dsi-ports-sortida-hexagonal.pptx"
LOGO = ASSETS_DIR / "logotip-oficial-tecnocampus-upf-horitzontal-color.png"
HEXASTOCK_ARCHITECTURE_IMAGE = ASSETS_DIR / "hexastock-sellstocks-arquitectura-vpd.png"

TITLE_COLOR = RGBColor(25, 45, 62)
TEXT_COLOR = RGBColor(38, 38, 38)
MUTED = RGBColor(102, 112, 122)
BLUE = RGBColor(35, 96, 146)
LIGHT_BLUE = RGBColor(232, 241, 248)
LIGHT_GREEN = RGBColor(235, 246, 240)
LIGHT_RED = RGBColor(248, 234, 234)
LIGHT_GREY = RGBColor(247, 248, 250)
YELLOW = RGBColor(255, 243, 214)


def set_slide_size(prs: Presentation) -> None:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_text(slide, text, x, y, w, h, font_size=22, bold=False, color=TEXT_COLOR,
             align=PP_ALIGN.LEFT, font_name="Arial"):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, font_size=20, color=TEXT_COLOR):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, item in enumerate(items):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_title(slide, title, subtitle=None):
    add_text(slide, title, Inches(0.65), Inches(0.35), Inches(9.8), Inches(0.55),
             font_size=28, bold=True, color=TITLE_COLOR)
    if subtitle:
        add_text(slide, subtitle, Inches(0.67), Inches(0.9), Inches(9.7), Inches(0.35),
                 font_size=13, color=MUTED)
    add_logo(slide, Inches(10.65), Inches(0.28), Inches(1.95), Inches(0.45))
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.22), Inches(12), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(214, 219, 224)
    line.line.fill.background()


def add_footer(slide, number):
    add_text(slide, "Disseny de Sistemes d'Informació · Ports de sortida en arquitectura hexagonal",
             Inches(0.65), Inches(7.08), Inches(9.4), Inches(0.22), font_size=8, color=MUTED)
    add_text(slide, str(number), Inches(12.2), Inches(7.08), Inches(0.45), Inches(0.22),
             font_size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def add_logo(slide, x, y, max_w, max_h):
    if not LOGO.exists():
        add_text(slide, "LOGOTIP OFICIAL TECNOCAMPUS / UPF", x, y, max_w, max_h,
                 font_size=8, color=MUTED, align=PP_ALIGN.RIGHT)
        return
    add_picture_fit(slide, LOGO, x, y, max_w, max_h)


def add_band(slide, x, y, w, h, fill, line=RGBColor(220, 224, 228)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_rect_band(slide, x, y, w, h, fill, line=RGBColor(220, 224, 228)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.75)
    return shape


def add_picture_fit(slide, path, x, y, max_w, max_h):
    path = Path(path)
    if not path.exists():
        placeholder = add_band(slide, x, y, max_w, max_h, LIGHT_GREY)
        add_text(slide, f"Pendent: {path.name}", x + Inches(0.15), y + Inches(0.15),
                 max_w - Inches(0.3), max_h - Inches(0.3), font_size=14, color=MUTED,
                 align=PP_ALIGN.CENTER)
        return placeholder

    with Image.open(path) as img:
        iw, ih = img.size
    ratio = min(max_w / iw, max_h / ih)
    width = int(iw * ratio)
    height = int(ih * ratio)
    left = x + int((max_w - width) / 2)
    top = y + int((max_h - height) / 2)
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_code(slide, code, x, y, w, h):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(31, 36, 41)
    shape.line.color.rgb = RGBColor(31, 36, 41)
    box = slide.shapes.add_textbox(x + Inches(0.22), y + Inches(0.18), w - Inches(0.44), h - Inches(0.36))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.name = "Consolas"
    run.font.size = Pt(15)
    run.font.color.rgb = RGBColor(245, 246, 247)
    return shape


def create_deck():
    prs = Presentation()
    set_slide_size(prs)

    # 1
    slide = blank_slide(prs)
    add_logo(slide, Inches(9.2), Inches(0.45), Inches(3.25), Inches(0.8))
    add_text(slide,
             "Disseny de ports de sortida en arquitectura hexagonal",
             Inches(0.72), Inches(1.35), Inches(10.8), Inches(0.9),
             font_size=32, bold=True, color=TITLE_COLOR)
    add_text(slide,
             "desacoblament entre casos d'ús, domini i APIs externes",
             Inches(0.75), Inches(2.18), Inches(10.2), Inches(0.45),
             font_size=21, color=BLUE)
    add_text(slide,
             "Una microlliçó de Disseny de Sistemes d'Informació a partir d'un cas real de consultoria i del projecte HexaStock",
             Inches(0.75), Inches(3.05), Inches(10.5), Inches(0.8),
             font_size=18, color=TEXT_COLOR)
    add_text(slide,
             "Alfredo Rueda Unsain\nDefensa de plaça de professorat permanent\nTecnoCampus · classe magistral reduïda, 20 minuts",
             Inches(0.75), Inches(5.1), Inches(7.5), Inches(0.85),
             font_size=15, color=MUTED)
    add_footer(slide, 1)

    # 2
    slide = blank_slide(prs)
    add_title(slide, "On som dins l'assignatura", "Disseny de Sistemes d'Informació")
    add_text(slide,
             "Disseny de Sistemes d'Informació\n3r curs · 6 ECTS",
             Inches(0.9), Inches(1.65), Inches(4.1), Inches(1.0), font_size=28, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_band(slide, Inches(5.25), Inches(1.65), Inches(0.05), Inches(3.5), BLUE, line=BLUE)
    add_text(slide,
             "Ports de sortida\nen arquitectura hexagonal",
             Inches(5.65), Inches(1.65), Inches(5.9), Inches(1.15), font_size=30, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "No explicarem tota l'arquitectura hexagonal.\n\nEns centrem en una decisió de disseny concreta.",
             Inches(5.9), Inches(3.25), Inches(5.45), Inches(1.25), font_size=22,
             color=TEXT_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Pla docent 103322: clean/hexagonal architecture, ports i adaptadors, mapping, mòduls i DDD.",
             Inches(1.05), Inches(6.2), Inches(10.8), Inches(0.35), font_size=13, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3
    slide = blank_slide(prs)
    add_title(slide, "Objectiu d'aprenentatge")
    add_text(slide,
             "El servei d'aplicació necessita una capacitat, no una tecnologia.",
             Inches(1.05), Inches(1.55), Inches(11.2), Inches(0.7), font_size=27, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    objectives = [
        "Identificar\nacoblament",
        "Definir\nun port",
        "Substituir\nadaptadors",
    ]
    for i, item in enumerate(objectives):
        x = Inches(0.9 + i * 4.08)
        add_band(slide, x, Inches(3.0), Inches(3.55), Inches(1.75), [LIGHT_RED, YELLOW, LIGHT_GREEN][i])
        add_text(slide, item, x + Inches(0.22), Inches(3.35), Inches(3.1), Inches(0.85),
                 font_size=24, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4
    slide = blank_slide(prs)
    add_title(slide, "Avaluació econòmica d'una beca",
              "Abans de parlar d'arquitectura, cal entendre el procediment")
    add_picture_fit(slide, DIAGRAMS_DIR / "agaur-functional-scholarship.png",
                    Inches(0.65), Inches(1.35), Inches(7.8), Inches(5.2))
    add_text(slide,
             "Sol·licitud → Expedient → Requisits generals\n→ Requisits econòmics → Revisió acadèmica\n→ Resolució",
             Inches(8.75), Inches(1.85), Inches(3.35), Inches(1.35), font_size=18,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_band(slide, Inches(8.75), Inches(3.65), Inches(3.35), Inches(1.35), LIGHT_RED)
    add_text(slide,
             "Si no es compleixen\nrequisits econòmics,\nl'expedient no avança.",
             Inches(9.0), Inches(3.9), Inches(2.85), Inches(0.82), font_size=20,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Fonts públiques: AGAUR requisits econòmics i BOE 2025-2026.",
             Inches(8.75), Inches(5.65), Inches(3.35), Inches(0.32), font_size=10, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5
    slide = blank_slide(prs)
    add_title(slide, "Quina informació externa necessita el procediment?",
              "Renda, patrimoni i informació cadastral")
    add_picture_fit(slide, DIAGRAMS_DIR / "agaur-data-interoperability.png",
                    Inches(0.65), Inches(1.35), Inches(7.8), Inches(5.15))
    add_text(slide, "NECESSITAT",
             Inches(8.95), Inches(1.55), Inches(3.2), Inches(0.35), font_size=16, bold=True,
             color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["Renda familiar", "Patrimoni", "Béns immobles"],
                Inches(9.05), Inches(2.0), Inches(3.0), Inches(1.05), font_size=17)
    add_text(slide, "FONTS",
             Inches(8.95), Inches(3.35), Inches(3.2), Inches(0.35), font_size=16, bold=True,
             color=BLUE, align=PP_ALIGN.CENTER)
    add_bullets(slide, ["AEAT", "Cadastre", "PICA"],
                Inches(9.05), Inches(3.8), Inches(3.0), Inches(1.0), font_size=17)
    add_text(slide,
             "El procediment necessita informació administrativa,\nno una API concreta.",
             Inches(8.85), Inches(5.35), Inches(3.35), Inches(0.75), font_size=16,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "PICA: interoperabilitat administrativa · Cadastre: registre de béns immobles · SOAP/XML: missatgeria i format d'integració",
             Inches(1.0), Inches(6.45), Inches(11.3), Inches(0.25), font_size=9, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # 6
    slide = blank_slide(prs)
    add_title(slide, "Quan canvia la integració, canvia massa codi")
    add_text(slide,
             "Cas d'ús  →  SOAP/XML  →  PICA  →  AEAT",
             Inches(1.05), Inches(1.75), Inches(11.2), Inches(0.85), font_size=31,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_rect_band(slide, Inches(1.0), Inches(3.0), Inches(3.35), Inches(1.3), LIGHT_RED)
    add_text(slide, "El cas d'ús coneix\nla infraestructura",
             Inches(1.25), Inches(3.3), Inches(2.85), Inches(0.6), font_size=18, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_rect_band(slide, Inches(4.95), Inches(3.0), Inches(3.35), Inches(1.3), YELLOW)
    add_text(slide, "El canvi tecnològic\nimpacta massa",
             Inches(5.2), Inches(3.3), Inches(2.85), Inches(0.6), font_size=18, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_rect_band(slide, Inches(8.9), Inches(3.0), Inches(3.35), Inches(1.3), LIGHT_BLUE)
    add_text(slide, "Errors tècnics afecten\nel procediment",
             Inches(9.15), Inches(3.3), Inches(2.85), Inches(0.6), font_size=18, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Separar casos d'ús, criteris del procediment i infraestructura.",
             Inches(1.05), Inches(5.55), Inches(11.2), Inches(0.55), font_size=25,
             bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7
    slide = blank_slide(prs)
    add_title(slide, "Port de sortida i adaptador")
    steps = [
        ("Cas d'ús", LIGHT_BLUE),
        ("Port de sortida", YELLOW),
        ("Adaptador", LIGHT_GREEN),
        ("Sistema extern", LIGHT_GREY),
    ]
    for i, (label, fill) in enumerate(steps):
        y = Inches(1.55 + i * 1.05)
        add_band(slide, Inches(4.15), y, Inches(4.9), Inches(0.62), fill)
        add_text(slide, label, Inches(4.35), y + Inches(0.11), Inches(4.5), Inches(0.28),
                 font_size=21, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text(slide, "↓", Inches(6.45), y + Inches(0.62), Inches(0.3), Inches(0.32),
                     font_size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "El port defineix què.",
             Inches(1.0), Inches(5.85), Inches(5.2), Inches(0.38), font_size=23,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "L'adaptador resol com.",
             Inches(7.0), Inches(5.85), Inches(5.2), Inches(0.38), font_size=23,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8
    slide = blank_slide(prs)
    add_title(slide, "De la interface al port", "Una frontera arquitectònica, no només un contracte Java")
    add_picture_fit(slide, DIAGRAMS_DIR / "output-port-pattern.png",
                    Inches(0.75), Inches(1.45), Inches(6.7), Inches(5.25))
    add_band(slide, Inches(8.45), Inches(1.75), Inches(3.65), Inches(3.8), LIGHT_GREY)
    add_text(slide,
             "Una interface\ndesacobla codi.\n\nUn port\ndesacobla arquitectura.",
             Inches(8.75), Inches(2.1), Inches(3.05), Inches(2.55), font_size=22, bold=True, color=TITLE_COLOR,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    # 9
    slide = blank_slide(prs)
    add_title(slide, "Transferència a HexaStock", "Venda d'accions amb proveïdor de preus substituïble")
    add_picture_fit(slide, HEXASTOCK_ARCHITECTURE_IMAGE,
                    Inches(0.55), Inches(1.35), Inches(8.6), Inches(5.55))
    add_text(slide,
             "Mateix problema,\ndiferent domini",
             Inches(9.35), Inches(1.62), Inches(3.0), Inches(0.8), font_size=22, bold=True,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "Cas d'ús: venda d'accions",
        "Necessitat: preu actual",
        "No depèn de proveïdors concrets",
    ], Inches(9.35), Inches(3.05), Inches(3.0), Inches(1.55), font_size=15)
    add_text(slide,
             "Port: StockPriceProviderPort\nAdaptadors: Finnhub, Alpha Vantage, mock",
             Inches(9.35), Inches(6.2), Inches(3.0), Inches(0.25), font_size=10, color=MUTED,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # 10
    slide = blank_slide(prs)
    add_title(slide, "Flux i codi essencial", "El servei coordina. El domini decideix. La infraestructura adapta.")
    add_picture_fit(slide, DIAGRAMS_DIR / "sell-stock-sequence.png",
                    Inches(0.55), Inches(1.35), Inches(8.0), Inches(4.95))
    code = (
        "stockPrice = port.fetch(ticker);\n\n"
        "result = portfolio.sell(...);"
    )
    add_code(slide, code, Inches(8.85), Inches(1.8), Inches(3.85), Inches(1.55))
    add_bullets(slide, [
        "El servei coordina",
        "El port obté informació",
        "El domini decideix",
    ], Inches(9.0), Inches(4.0), Inches(3.35), Inches(1.25), font_size=16)
    add_footer(slide, 10)

    # 11
    slide = blank_slide(prs)
    add_title(slide, "Demo controlada", "Mateix cas d'ús, adaptador substituïble")
    add_text(slide,
             "HexaStock és un projecte open source propi utilitzat en docència universitària i formació o consultoria amb institucions financeres.",
             Inches(1.0), Inches(1.45), Inches(11.25), Inches(0.55), font_size=18,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    for i, label in enumerate(["Mateix cas d'ús", "Mateix servei", "Mateix domini"]):
        x = Inches(0.95 + i * 4.05)
        add_band(slide, x, Inches(2.45), Inches(3.55), Inches(1.05), [LIGHT_BLUE, YELLOW, LIGHT_GREEN][i])
        add_text(slide, label, x + Inches(0.2), Inches(2.78), Inches(3.15), Inches(0.28),
                 font_size=19, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Canvi: adaptador real  ↔  adaptador mock",
             Inches(1.15), Inches(4.25), Inches(11.0), Inches(0.55), font_size=26,
             bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Canvia la infraestructura; el cas d'ús roman estable.",
             Inches(1.15), Inches(5.55), Inches(11.0), Inches(0.5), font_size=24,
             bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_footer(slide, 11)

    # 12
    slide = blank_slide(prs)
    add_title(slide, "Conclusió", "L'arquitectura no elimina el canvi. El localitza.")
    add_picture_fit(slide, DIAGRAMS_DIR / "before-after-comparison.png",
                    Inches(0.75), Inches(1.45), Inches(7.75), Inches(5.05))
    add_text(slide,
             "Quan el proveïdor canvia, volem canviar l'adaptador, no el cas d'ús ni el model de domini.",
             Inches(8.95), Inches(2.05), Inches(3.25), Inches(1.9), font_size=24, bold=True, color=TITLE_COLOR,
             align=PP_ALIGN.CENTER)
    add_footer(slide, 12)

    # 13
    slide = blank_slide(prs)
    add_title(slide, "Agraïment")
    add_text(slide,
             "Vull expressar el meu agraïment al TecnoCampus per l'oportunitat d'aprendre, créixer com a docent i transferir coneixement entre l'empresa i la universitat.",
             Inches(1.35), Inches(1.75), Inches(10.5), Inches(1.1), font_size=24,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_text(slide,
             "Aquesta experiència m'ha permès portar casos reals, criteri professional i arquitectura aplicada a l'aula.",
             Inches(1.65), Inches(3.25), Inches(9.9), Inches(0.75), font_size=21,
             color=TEXT_COLOR, align=PP_ALIGN.CENTER)
    add_band(slide, Inches(2.15), Inches(4.75), Inches(9.05), Inches(1.0), LIGHT_GREY)
    add_text(slide,
             "Agraïment especial al Dr. Josep Roure, company professor del TecnoCampus, amb qui he tingut l'oportunitat d'aprendre molt sobre arquitectura i software.",
             Inches(2.45), Inches(4.98), Inches(8.45), Inches(0.48), font_size=16,
             color=TITLE_COLOR, align=PP_ALIGN.CENTER)
    add_footer(slide, 13)

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"PowerPoint generat: {OUTPUT_FILE}")


if __name__ == "__main__":
    create_deck()
